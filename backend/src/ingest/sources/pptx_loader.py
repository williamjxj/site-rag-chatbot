"""PowerPoint document loader (.pptx)."""

from pathlib import Path
from pptx import Presentation


def load_pptx(path: Path) -> dict[str, str]:
    """
    Load and extract text from PowerPoint .pptx file.

    Args:
        path: Path to .pptx file

    Returns:
        Dictionary with 'uri', 'title', and 'text' keys
    """
    prs = Presentation(str(path))
    parts = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        parts.append(f"Slide {slide_num}")
        
        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                # Add shape type as context if available
                if hasattr(shape, "shape_type"):
                    parts.append(f"  [{shape.shape_type}]: {text}")
                else:
                    parts.append(f"  {text}")
        
        parts.append("")  # Empty line between slides
    
    text = "\n".join(parts).strip()
    return {"uri": str(path), "title": path.stem, "text": text}

